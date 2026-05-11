"""
html_to_pptx.py - Export the adversarial-patch-only presentation deck to PowerPoint.

Usage:
    python scripts/html_to_pptx.py
"""

from pathlib import Path
from shutil import copy2

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BG = RGBColor(0x0A, 0x0A, 0x0F)
WHITE = RGBColor(0xE8, 0xE8, 0xF0)
DIM = RGBColor(0x88, 0x88, 0x88)
DIMMER = RGBColor(0x44, 0x44, 0x44)
BODY = RGBColor(0xCC, 0xCC, 0xCC)
GREEN = RGBColor(0x3E, 0xFF, 0xA0)
YELLOW = RGBColor(0xFF, 0xE0, 0x66)
RED = RGBColor(0xFF, 0x60, 0x60)
BLUE = RGBColor(0x60, 0xB8, 0xFF)
CARD_BG = RGBColor(0x12, 0x12, 0x1A)
CARD_HL = RGBColor(0x0D, 0x1A, 0x12)
CARD_WN = RGBColor(0x1A, 0x0D, 0x0D)
CARD_INF = RGBColor(0x0D, 0x12, 0x20)
CARD_GOLD = RGBColor(0x1A, 0x18, 0x0D)

W = Inches(13.333)
H = Inches(7.5)
ROOT = Path(__file__).parent.parent
ASSET_DIR = ROOT / "docs" / "assets" / "presentation"


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide


def rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    return shape


def textbox(
    slide,
    text,
    left,
    top,
    width,
    height,
    *,
    size=18,
    bold=False,
    color=WHITE,
    align=PP_ALIGN.LEFT,
    italic=False,
):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def rich_textbox(slide, lines, left, top, width, height):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for spec in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = spec.get("align", PP_ALIGN.LEFT)
        run = p.add_run()
        run.text = spec["text"]
        run.font.size = Pt(spec.get("size", 18))
        run.font.bold = spec.get("bold", False)
        run.font.italic = spec.get("italic", False)
        run.font.color.rgb = spec.get("color", WHITE)
    return txb


def add_contained_picture(slide, image_path: Path, left, top, width, height):
    if not image_path.exists():
        return None
    with Image.open(image_path) as img:
        img_w, img_h = img.size
    box_ratio = float(width) / float(height)
    img_ratio = img_w / img_h
    if img_ratio > box_ratio:
        pic_w = width
        pic_h = int(float(width) / img_ratio)
        pic_left = left
        pic_top = top + int((float(height) - pic_h) / 2)
    else:
        pic_h = height
        pic_w = int(float(height) * img_ratio)
        pic_left = left + int((float(width) - pic_w) / 2)
        pic_top = top
    return slide.shapes.add_picture(str(image_path), pic_left, pic_top, pic_w, pic_h)


def picture_card(slide, image_path: Path, left, top, width, height, *, fill=CARD_BG, inner_fill=BG):
    rect(slide, left, top, width, height, fill)
    inner_pad = Inches(0.14)
    rect(slide, left + inner_pad, top + inner_pad, width - inner_pad * 2, height - inner_pad * 2, inner_fill)
    return add_contained_picture(
        slide,
        image_path,
        left + inner_pad,
        top + inner_pad,
        width - inner_pad * 2,
        height - inner_pad * 2,
    )


def rule(slide, left, top, width, color=RGBColor(0x1E, 0x1E, 0x28)):
    rect(slide, left, top, width, Inches(0.02), color)


def tag(slide, text, left=Inches(0.9), top=Inches(0.45)):
    textbox(slide, text, left, top, Inches(4.8), Inches(0.28), size=9, bold=True, color=BLUE)


def title_block(slide, lines, top=Inches(0.85), left=Inches(0.9)):
    rich_textbox(slide, lines, left, top, Inches(11.5), Inches(1.7))


def card_title(slide, text, left, top, width):
    textbox(slide, text, left, top, width, Inches(0.22), size=9, bold=True, color=BLUE)


def bullet_lines(items, color=BODY, size=12):
    return [{"text": f"- {item}", "size": size, "color": color} for item in items]


def slide1_title(prs: Presentation):
    slide = blank_slide(prs)
    tag(slide, "AI FOR CYBERSECURITY FINAL PROJECT", top=Inches(0.68))
    title_block(
        slide,
        [
            {"text": "Adversarial Patch Attacks on", "size": 34, "bold": True, "color": WHITE},
            {"text": "YOLO Person Detection", "size": 34, "bold": True, "color": GREEN},
        ],
        top=Inches(1.1),
    )
    textbox(
        slide,
        "Stand-alone presentation for the Adversarial_Patch research pipeline",
        Inches(0.9),
        Inches(2.72),
        Inches(7.0),
        Inches(0.42),
        size=15,
        color=DIM,
    )
    textbox(
        slide,
        "This deck focuses only on the adversarial patch project: how a small localized patch was trained, evaluated, and analyzed against YOLO person detection models.",
        Inches(0.9),
        Inches(3.45),
        Inches(5.45),
        Inches(1.05),
        size=14,
        color=BODY,
    )
    rule(slide, Inches(0.9), Inches(4.8), Inches(1.2), RGBColor(0x3E, 0xFF, 0xA0))
    textbox(slide, "April 2026", Inches(0.9), Inches(5.02), Inches(2.2), Inches(0.35), size=12, color=DIMMER)

    left = Inches(6.85)
    top = Inches(2.2)
    width = Inches(5.55)
    height = Inches(3.9)
    rect(slide, left, top, width, height, CARD_INF)
    card_title(slide, "TITLE & META-DATA", left + Inches(0.2), top + Inches(0.16), width)
    rich_textbox(
        slide,
        [
            {"text": "Project Title: Adversarial Patch Attacks on YOLO Person Detection", "size": 12, "color": BODY},
            {"text": "Name(s): [Your Name(s)]", "size": 12, "color": BODY},
            {"text": "Category: Adversarial ML / Computer Vision Security", "size": 12, "color": BODY},
            {
                "text": "Capstone relation: Smaller AI-for-Cybersecurity project derived from the broader course capstone; this presentation covers only the adversarial patch subproject.",
                "size": 12,
                "color": BODY,
            },
        ],
        left + Inches(0.2),
        top + Inches(0.55),
        width - Inches(0.35),
        Inches(3.0),
    )


def slide2_quad_chart(prs: Presentation, patch_path: Path):
    slide = blank_slide(prs)
    tag(slide, "QUAD CHART")
    title_block(
        slide,
        [
            {"text": "Project at a", "size": 36, "bold": True, "color": WHITE},
            {"text": "Glance", "size": 36, "bold": True, "color": BLUE},
        ],
    )

    x1 = Inches(0.9)
    x2 = Inches(6.72)
    y1 = Inches(2.0)
    y2 = Inches(4.1)
    w = Inches(5.75)
    h = Inches(1.72)

    rect(slide, x1, y1, w, h, CARD_INF)
    card_title(slide, "PROBLEM", x1 + Inches(0.2), y1 + Inches(0.16), w)
    textbox(
        slide,
        "YOLO-based person detection can be degraded by a small localized patch. That matters anywhere automated vision is trusted for awareness, safety, or monitoring.",
        x1 + Inches(0.2),
        y1 + Inches(0.5),
        w - Inches(0.35),
        Inches(0.95),
        size=12,
        color=BODY,
    )

    rect(slide, x2, y1, w, h, CARD_BG)
    card_title(slide, "APPROACH", x2 + Inches(0.2), y1 + Inches(0.16), w)
    textbox(
        slide,
        "Train and test a 100 x 100 adversarial patch across YOLOv8n, YOLO11n, and YOLO26n, then measure direct suppression, transfer, and follow-up variants.",
        x2 + Inches(0.2),
        y1 + Inches(0.5),
        w - Inches(0.35),
        Inches(0.95),
        size=12,
        color=BODY,
    )

    rect(slide, x1, y2, w, h, CARD_HL)
    card_title(slide, "RESULTS", x1 + Inches(0.2), y2 + Inches(0.16), w)
    textbox(
        slide,
        "Direct suppression reached 90.0% on YOLOv8n and 72.7% on YOLO11n, while YOLO26n resisted with only 16.3%. Transfer existed, but it was asymmetric.",
        x1 + Inches(0.2),
        y2 + Inches(0.5),
        w - Inches(0.35),
        Inches(0.95),
        size=12,
        color=BODY,
    )

    rect(slide, x2, y2, w, h, CARD_WN)
    card_title(slide, "IMPACT", x2 + Inches(0.2), y2 + Inches(0.16), w)
    rect(slide, x2 + Inches(0.22), y2 + Inches(0.48), Inches(0.9), Inches(0.9), CARD_BG)
    add_contained_picture(slide, patch_path, x2 + Inches(0.29), y2 + Inches(0.55), Inches(0.76), Inches(0.76))
    textbox(
        slide,
        "The project produced a learned patch artifact, proof images, transfer analysis, and print-ready outputs that make the attack concrete rather than hypothetical.",
        x2 + Inches(1.25),
        y2 + Inches(0.48),
        Inches(4.2),
        Inches(0.95),
        size=12,
        color=BODY,
    )


def slide3_abstract(prs: Presentation):
    slide = blank_slide(prs)
    tag(slide, "ABSTRACT")
    title_block(
        slide,
        [
            {"text": "Brief", "size": 36, "bold": True, "color": WHITE},
            {"text": "Overview", "size": 36, "bold": True, "color": GREEN},
        ],
    )

    rect(slide, Inches(0.9), Inches(2.0), Inches(7.2), Inches(4.7), CARD_HL)
    card_title(slide, "ABSTRACT", Inches(1.1), Inches(2.16), Inches(6.8))
    textbox(
        slide,
        "This project tested whether a small localized adversarial patch could suppress person detection across modern Ultralytics YOLO models. A 100 x 100 patch was trained and evaluated against YOLOv8n, YOLO11n, and YOLO26n using direct, transfer, joint, and warm-start experiments. Results showed strong direct suppression on v8n and v11n, measurable but asymmetric transfer, and a major architectural limit on YOLO26n. The project also produced reusable artifacts including learned patches, proof images, transfer visualizations, and print-ready outputs. Together these results show both the offensive potential of adversarial patches and the model-specific barriers that matter for real attack development.",
        Inches(1.1),
        Inches(2.55),
        Inches(6.8),
        Inches(3.55),
        size=14,
        color=BODY,
    )

    right = Inches(8.45)
    rect(slide, right, Inches(2.0), Inches(3.95), Inches(1.35), CARD_INF)
    card_title(slide, "SCOPE", right + Inches(0.2), Inches(2.16), Inches(3.6))
    textbox(
        slide,
        "Models studied: YOLOv8n, YOLO11n, and YOLO26n. Focus: person detection suppression through a small visible patch.",
        right + Inches(0.2),
        Inches(2.48),
        Inches(3.55),
        Inches(0.7),
        size=11,
        color=BODY,
    )

    rect(slide, right, Inches(3.58), Inches(3.95), Inches(1.25), CARD_BG)
    card_title(slide, "CORE FINDING", right + Inches(0.2), Inches(3.74), Inches(3.6))
    textbox(
        slide,
        "The patch is clearly effective on two model generations, but not all YOLO architectures fail in the same way.",
        right + Inches(0.2),
        Inches(4.05),
        Inches(3.55),
        Inches(0.6),
        size=11,
        color=BODY,
    )

    rect(slide, right, Inches(5.03), Inches(3.95), Inches(1.65), CARD_WN)
    card_title(slide, "WHY IT MATTERS", right + Inches(0.2), Inches(5.19), Inches(3.6))
    textbox(
        slide,
        "Attack research is more useful when it produces evidence, artifacts, and limitations that can be demonstrated and discussed directly.",
        right + Inches(0.2),
        Inches(5.5),
        Inches(3.55),
        Inches(0.9),
        size=11,
        color=BODY,
    )


def slide4_introduction(prs: Presentation):
    slide = blank_slide(prs)
    tag(slide, "INTRODUCTION")
    title_block(
        slide,
        [
            {"text": "Objective and", "size": 36, "bold": True, "color": WHITE},
            {"text": "Ultimate Impact", "size": 36, "bold": True, "color": GREEN},
        ],
    )

    rect(slide, Inches(0.9), Inches(2.0), Inches(5.75), Inches(2.25), CARD_INF)
    card_title(slide, "OBJECTIVE", Inches(1.1), Inches(2.16), Inches(5.3))
    textbox(
        slide,
        "Determine whether a small portable patch can reliably suppress person detection across current YOLO generations instead of only a single detector.",
        Inches(1.1),
        Inches(2.52),
        Inches(5.35),
        Inches(0.82),
        size=13,
        color=BODY,
    )
    textbox(
        slide,
        "Success meant demonstrating measurable suppression, cross-model behavior, and an evidence-based explanation when the attack failed.",
        Inches(1.1),
        Inches(3.42),
        Inches(5.35),
        Inches(0.48),
        size=11,
        color=BODY,
    )

    rect(slide, Inches(6.72), Inches(2.0), Inches(5.7), Inches(2.25), CARD_HL)
    card_title(slide, "ULTIMATE IMPACT", Inches(6.92), Inches(2.16), Inches(5.3))
    textbox(
        slide,
        "If a simple patch can degrade person detection, then surveillance, robotics, and automated vision systems may be vulnerable to localized physical or digital perturbations that are cheap to deploy.",
        Inches(6.92),
        Inches(2.52),
        Inches(5.3),
        Inches(1.15),
        size=13,
        color=BODY,
    )

    rect(slide, Inches(0.9), Inches(4.65), Inches(11.52), Inches(1.65), CARD_WN)
    card_title(slide, "THREAT MODEL", Inches(1.1), Inches(4.82), Inches(11.0))
    textbox(
        slide,
        "An attacker introduces a visible patch on clothing or another surface near the target person. The rest of the scene remains unchanged; the goal is to reduce or remove person detections rather than corrupt the full image.",
        Inches(1.1),
        Inches(5.15),
        Inches(11.1),
        Inches(0.8),
        size=12,
        color=BODY,
    )


def slide5_background(prs: Presentation):
    slide = blank_slide(prs)
    tag(slide, "BACKGROUND")
    title_block(
        slide,
        [
            {"text": "Problem Space and", "size": 34, "bold": True, "color": WHITE},
            {"text": "Limits of Current Practice", "size": 34, "bold": True, "color": BLUE},
        ],
    )

    rect(slide, Inches(0.9), Inches(2.0), Inches(5.75), Inches(2.55), CARD_BG)
    card_title(slide, "PROBLEM SPACE", Inches(1.1), Inches(2.16), Inches(5.3))
    textbox(
        slide,
        "Adversarial machine learning has shown that small perturbations can mislead vision models. For object detection, the challenge is harder than image classification because the attack must alter localization, confidence, and detection persistence in realistic scenes.",
        Inches(1.1),
        Inches(2.52),
        Inches(5.35),
        Inches(1.55),
        size=12,
        color=BODY,
    )

    rect(slide, Inches(6.72), Inches(2.0), Inches(5.7), Inches(2.55), CARD_INF)
    card_title(slide, "LIMITS OF CURRENT PRACTICE", Inches(6.92), Inches(2.16), Inches(5.3))
    rich_textbox(
        slide,
        bullet_lines(
            [
                "Many studies stop at one detector or one model family.",
                "Older-model results do not guarantee the same failure mode on newer architectures.",
                "Simple preprocessing defenses are often assumed to help without strong cross-model validation.",
            ],
            BODY,
            11,
        ),
        Inches(6.92),
        Inches(2.52),
        Inches(5.25),
        Inches(1.55),
    )

    rect(slide, Inches(0.9), Inches(4.92), Inches(11.52), Inches(1.4), CARD_HL)
    card_title(slide, "GAP ADDRESSED BY THIS PROJECT", Inches(1.1), Inches(5.08), Inches(11.0))
    textbox(
        slide,
        "This work compares direct suppression, transfer, and follow-up attack variants across three YOLO generations so the conclusions are tied to evidence rather than to a single-model anecdote.",
        Inches(1.1),
        Inches(5.38),
        Inches(11.1),
        Inches(0.7),
        size=12,
        color=BODY,
    )


def slide6_methodology(prs: Presentation, method_path: Path):
    slide = blank_slide(prs)
    tag(slide, "METHODOLOGY")
    title_block(
        slide,
        [
            {"text": "Novel", "size": 36, "bold": True, "color": WHITE},
            {"text": "Approach", "size": 36, "bold": True, "color": GREEN},
        ],
    )

    picture_card(slide, method_path, Inches(0.9), Inches(2.0), Inches(6.1), Inches(3.55), fill=CARD_BG, inner_fill=BG)
    textbox(
        slide,
        "Train the patch, overlay it on the target person, evaluate detections, then measure transfer and follow-up variants.",
        Inches(1.05),
        Inches(5.62),
        Inches(5.8),
        Inches(0.35),
        size=10,
        color=DIM,
    )

    right = Inches(7.25)
    rect(slide, right, Inches(2.0), Inches(5.15), Inches(1.7), CARD_INF)
    card_title(slide, "EXPERIMENTAL SETUP", right + Inches(0.2), Inches(2.16), Inches(4.8))
    rich_textbox(
        slide,
        bullet_lines(
            [
                "Patch size: 100 x 100 pixels",
                "Image size: 640 x 640",
                "Common 48-image manifest for cross-model consistency",
                "Evaluation on each model's clean person-detection subset",
            ],
            BODY,
            11,
        ),
        right + Inches(0.2),
        Inches(2.48),
        Inches(4.75),
        Inches(1.15),
    )

    rect(slide, right, Inches(3.95), Inches(5.15), Inches(1.6), CARD_BG)
    card_title(slide, "NOVEL ELEMENTS", right + Inches(0.2), Inches(4.11), Inches(4.8))
    rich_textbox(
        slide,
        bullet_lines(
            [
                "Direct self-evaluation on each YOLO model",
                "Cross-model transfer matrix across generations",
                "Joint-patch and warm-start follow-up experiments",
                "Exportable proof and print artifacts for demonstration",
            ],
            BODY,
            10,
        ),
        right + Inches(0.2),
        Inches(4.42),
        Inches(4.75),
        Inches(1.05),
    )

    rect(slide, Inches(0.9), Inches(6.0), Inches(11.52), Inches(0.92), CARD_WN)
    card_title(slide, "HOW SUCCESS WAS PROVEN", Inches(1.1), Inches(6.15), Inches(11.0))
    textbox(
        slide,
        "Success was proven with clean-versus-patched person detections, suppression percentage, transfer outcomes, and follow-up experiments that tested whether weak results came from initialization or architecture.",
        Inches(1.1),
        Inches(6.36),
        Inches(11.1),
        Inches(0.35),
        size=10,
        color=BODY,
    )


def slide7_direct_results(prs: Presentation, proof_path: Path):
    slide = blank_slide(prs)
    tag(slide, "RESULTS & DISCUSSION")
    title_block(
        slide,
        [
            {"text": "Direct", "size": 36, "bold": True, "color": WHITE},
            {"text": "Suppression", "size": 36, "bold": True, "color": GREEN},
        ],
    )

    picture_card(slide, proof_path, Inches(0.9), Inches(2.0), Inches(7.15), Inches(4.95), fill=CARD_BG, inner_fill=CARD_INF)
    textbox(
        slide,
        "Primary proof object: the same scene before and after the patch, with the clean person detection removed in the patched frame.",
        Inches(1.05),
        Inches(6.58),
        Inches(6.85),
        Inches(0.35),
        size=10,
        color=DIM,
    )

    card_x = Inches(8.35)
    card_w = Inches(4.05)
    card_h = Inches(1.05)
    starts = [Inches(2.08), Inches(3.32), Inches(4.56)]
    values = [
        ("90.0%", "YOLOv8n", "20 -> 2 detections | strongest direct result", CARD_HL, GREEN),
        ("72.7%", "YOLO11n", "33 -> 9 detections | strong newer-model suppression", CARD_GOLD, YELLOW),
        ("16.3%", "YOLO26n", "43 -> 36 detections | weak suppression despite low loss", CARD_WN, RED),
    ]
    for top, (value, title, detail, fill, color) in zip(starts, values):
        rect(slide, card_x, top, card_w, card_h, fill)
        textbox(slide, value, card_x + Inches(0.2), top + Inches(0.16), Inches(1.55), Inches(0.45), size=29, bold=True, color=color)
        textbox(slide, title, card_x + Inches(1.82), top + Inches(0.16), Inches(1.95), Inches(0.22), size=13, bold=True)
        textbox(slide, detail, card_x + Inches(1.82), top + Inches(0.48), Inches(1.95), Inches(0.28), size=9, color=DIM)

    rect(slide, card_x, Inches(5.84), card_w, Inches(1.1), CARD_INF)
    card_title(slide, "DISCUSSION", card_x + Inches(0.2), Inches(5.98), card_w)
    textbox(
        slide,
        "v8n and v11n are strongly suppressible. YOLO26n is the exception, which makes the study more useful than a single success case.",
        card_x + Inches(0.2),
        Inches(6.2),
        Inches(3.6),
        Inches(0.48),
        size=10,
        color=BODY,
    )


def slide8_transfer(prs: Presentation, heatmap_path: Path):
    slide = blank_slide(prs)
    tag(slide, "RESULTS & DISCUSSION")
    title_block(
        slide,
        [
            {"text": "Transfer Across", "size": 34, "bold": True, "color": WHITE},
            {"text": "Models", "size": 34, "bold": True, "color": BLUE},
        ],
    )

    picture_card(slide, heatmap_path, Inches(0.9), Inches(2.0), Inches(6.95), Inches(4.95), fill=CARD_BG, inner_fill=BG)
    textbox(
        slide,
        "Transfer exists across model generations, but it is not symmetric and it is weakest when YOLO26n is the target.",
        Inches(1.05),
        Inches(6.58),
        Inches(6.65),
        Inches(0.35),
        size=10,
        color=DIM,
    )

    right = Inches(8.12)
    rect(slide, right, Inches(2.0), Inches(4.3), Inches(1.2), CARD_HL)
    card_title(slide, "ASYMMETRY", right + Inches(0.2), Inches(2.15), Inches(4.0))
    textbox(
        slide,
        "v11n -> v8n reached 55.0%, while v8n -> v11n reached 33.3%. Adversarial features did not transfer equally in both directions.",
        right + Inches(0.2),
        Inches(2.48),
        Inches(3.95),
        Inches(0.62),
        size=11,
        color=BODY,
    )

    rect(slide, right, Inches(3.45), Inches(4.3), Inches(1.7), CARD_INF)
    card_title(slide, "TRANSFER MATRIX", right + Inches(0.2), Inches(3.6), Inches(4.0))
    rich_textbox(
        slide,
        [
            {"text": "v8n -> v11n: 33.3%", "size": 10, "color": BODY},
            {"text": "v11n -> v8n: 55.0%", "size": 10, "color": BODY},
            {"text": "v26n -> v8n: 45.0%", "size": 10, "color": BODY},
            {"text": "v26n -> v11n: 24.2%", "size": 10, "color": BODY},
            {"text": "v8n -> v26n: 14.0%", "size": 10, "color": BODY},
            {"text": "v11n -> v26n: 9.3%", "size": 10, "color": BODY},
        ],
        right + Inches(0.2),
        Inches(3.92),
        Inches(3.95),
        Inches(1.2),
    )

    rect(slide, right, Inches(5.42), Inches(4.3), Inches(1.53), CARD_BG)
    card_title(slide, "FOLLOW-UP JOINT RESULT", right + Inches(0.2), Inches(5.57), Inches(4.0))
    textbox(
        slide,
        "Even the best v26n-targeting joint patch only reached 18.6%, so the barrier is not just single-model overfitting.",
        right + Inches(0.2),
        Inches(5.88),
        Inches(3.95),
        Inches(0.7),
        size=10,
        color=BODY,
    )


def slide9_yolo26n(prs: Presentation, paradox_path: Path, print_path: Path):
    slide = blank_slide(prs)
    tag(slide, "RESULTS & DISCUSSION")
    title_block(
        slide,
        [
            {"text": "YOLO26n Limitation and", "size": 32, "bold": True, "color": WHITE},
            {"text": "Tool Outputs", "size": 32, "bold": True, "color": RED},
        ],
    )

    picture_card(slide, paradox_path, Inches(0.9), Inches(2.0), Inches(6.0), Inches(4.95), fill=CARD_BG, inner_fill=BG)
    textbox(
        slide,
        "Low loss did not translate into strong suppression, which made YOLO26n the project's most important architecture-specific finding.",
        Inches(1.05),
        Inches(6.58),
        Inches(5.7),
        Inches(0.35),
        size=10,
        color=DIM,
    )

    right = Inches(7.15)
    rect(slide, right, Inches(2.0), Inches(5.25), Inches(1.35), CARD_WN)
    card_title(slide, "ARCHITECTURAL FINDING", right + Inches(0.2), Inches(2.15), Inches(4.8))
    textbox(
        slide,
        "Direct suppression stalled at 16.3% despite low final detection loss, likely because training and inference use different detection heads.",
        right + Inches(0.2),
        Inches(2.48),
        Inches(4.9),
        Inches(0.85),
        size=10,
        color=BODY,
    )

    rect(slide, right, Inches(3.55), Inches(5.25), Inches(1.2), CARD_INF)
    card_title(slide, "FOLLOW-UP CHECKS", right + Inches(0.2), Inches(3.7), Inches(4.8))
    textbox(
        slide,
        "Warm-start reached 14.0%. A one2one follow-up reached 11.6%. That points to a structural limit, not a bad initialization.",
        right + Inches(0.2),
        Inches(4.02),
        Inches(4.9),
        Inches(0.62),
        size=10,
        color=BODY,
    )

    rect(slide, right, Inches(4.95), Inches(3.2), Inches(1.98), CARD_BG)
    card_title(slide, "TOOL OUTPUTS", right + Inches(0.2), Inches(5.1), Inches(2.8))
    rich_textbox(
        slide,
        [
            {"text": "- learned patch artifact", "size": 9, "color": BODY},
            {"text": "- clean-versus-patched proof images", "size": 9, "color": BODY},
            {"text": "- transfer heatmap", "size": 9, "color": BODY},
            {"text": "- print-ready export", "size": 9, "color": BODY},
            {"text": "- JPEG/blur usually failed; crop-resize had isolated wins", "size": 9, "color": BODY},
        ],
        right + Inches(0.2),
        Inches(5.35),
        Inches(2.75),
        Inches(1.4),
    )

    picture_card(slide, print_path, Inches(10.62), Inches(5.0), Inches(1.78), Inches(1.35), fill=CARD_INF, inner_fill=BG)


def slide10_conclusion(prs: Presentation):
    slide = blank_slide(prs)
    tag(slide, "CONCLUSION")
    title_block(
        slide,
        [
            {"text": "Takeaways, Limits, and", "size": 32, "bold": True, "color": WHITE},
            {"text": "Future Work", "size": 32, "bold": True, "color": GREEN},
        ],
    )

    card_w = Inches(3.6)
    card_h = Inches(2.05)
    y = Inches(2.1)
    gap = Inches(0.36)
    x1 = Inches(0.9)
    x2 = x1 + card_w + gap
    x3 = x2 + card_w + gap

    rect(slide, x1, y, card_w, card_h, CARD_HL)
    card_title(slide, "TAKEAWAY 1", x1 + Inches(0.2), y + Inches(0.16), card_w)
    textbox(
        slide,
        "Localized adversarial patches can strongly suppress person detection on some YOLO generations, especially YOLOv8n and YOLO11n.",
        x1 + Inches(0.2),
        y + Inches(0.5),
        card_w - Inches(0.35),
        Inches(1.2),
        size=12,
        color=BODY,
    )

    rect(slide, x2, y, card_w, card_h, CARD_INF)
    card_title(slide, "TAKEAWAY 2", x2 + Inches(0.2), y + Inches(0.16), card_w)
    textbox(
        slide,
        "Transfer across models is real but asymmetric, so attack behavior cannot be summarized by a single success rate.",
        x2 + Inches(0.2),
        y + Inches(0.5),
        card_w - Inches(0.35),
        Inches(1.2),
        size=12,
        color=BODY,
    )

    rect(slide, x3, y, card_w, card_h, CARD_WN)
    card_title(slide, "MAIN LIMITATION", x3 + Inches(0.2), y + Inches(0.16), card_w)
    textbox(
        slide,
        "YOLO26n resisted the current attack formulation, which means better optimization alone is unlikely to solve the problem.",
        x3 + Inches(0.2),
        y + Inches(0.5),
        card_w - Inches(0.35),
        Inches(1.2),
        size=12,
        color=BODY,
    )

    rect(slide, Inches(0.9), Inches(4.55), Inches(11.52), Inches(1.55), CARD_BG)
    card_title(slide, "FUTURE WORK", Inches(1.1), Inches(4.72), Inches(11.0))
    textbox(
        slide,
        "Next steps are to design a matching-aware objective for YOLO26n, expand physical benchmarking, study placement robustness, and evaluate whether stronger cross-model attacks can preserve suppression without sacrificing portability.",
        Inches(1.1),
        Inches(5.05),
        Inches(11.1),
        Inches(0.8),
        size=12,
        color=BODY,
    )


def main():
    patch_path = ROOT / "outputs" / "yolov8n_patch_v2" / "patches" / "patch.png"
    proof_path = ASSET_DIR / "yolov8n_proof_pair.png"
    method_path = ASSET_DIR / "method_pipeline.png"
    heatmap_path = ASSET_DIR / "transfer_heatmap.png"
    paradox_path = ASSET_DIR / "yolo26n_paradox.png"
    print_path = ASSET_DIR / "patch_print_inset.png"
    final_path = ROOT / "adversarial_patch_presentation.pptx"
    legacy_path = ROOT / "deck.pptx"

    prs = new_prs()
    slide1_title(prs)
    slide2_quad_chart(prs, patch_path)
    slide3_abstract(prs)
    slide4_introduction(prs)
    slide5_background(prs)
    slide6_methodology(prs, method_path)
    slide7_direct_results(prs, proof_path)
    slide8_transfer(prs, heatmap_path)
    slide9_yolo26n(prs, paradox_path, print_path)
    slide10_conclusion(prs)
    prs.save(final_path)
    copy2(final_path, legacy_path)
    print(f"Saved: {final_path}")
    print(f"Copied: {legacy_path}")


if __name__ == "__main__":
    main()
